import os
import re
import json
import hashlib
import logging
import secrets
import warnings
from datetime import datetime, timedelta
from functools import cached_property, wraps
from pathlib import Path
from collections import defaultdict, OrderedDict
from contextlib import contextmanager

import pandas as pd
import numpy as np
import requests
from flask import Flask, jsonify, request, render_template, send_file
from flask_cors import CORS
from flask_compress import Compress
from werkzeug.middleware.proxy_fix import ProxyFix
import jwt
import psycopg2
from psycopg2.pool import SimpleConnectionPool
from psycopg2.extras import RealDictCursor
from cachetools import TTLCache
from sklearn.linear_model import LinearRegression
import io
import tempfile
import sys
from dotenv import load_dotenv

load_dotenv()
warnings.filterwarnings('ignore')

# ------------------------------
#  Logging setup
# ------------------------------
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# ------------------------------
#  Configuration (environment based)
# ------------------------------
class Config:
    SECRET_KEY = os.environ.get('SECRET_KEY', os.urandom(32))
    SESSION_COOKIE_SECURE = True
    SESSION_COOKIE_HTTPONLY = True
    SESSION_COOKIE_SAMESITE = 'Lax'
    PERMANENT_SESSION_LIFETIME = timedelta(hours=1)
    
    RATELIMIT_AI = int(os.environ.get('RATELIMIT_AI', 20))
    RATELIMIT_NLQ = int(os.environ.get('RATELIMIT_NLQ', 50))
    RATELIMIT_SIMULATE = int(os.environ.get('RATELIMIT_SIMULATE', 30))
    
    DATABASE_URL = os.environ.get('DATABASE_URL')
    if not DATABASE_URL:
        print("\n❌ ERROR: DATABASE_URL environment variable not set.")
        print("   Please create a .env file with your Neon PostgreSQL URL.")
        print("   Example .env content:")
        print("   DATABASE_URL=postgresql://user:pass@host:5432/db")
        print("   Then run: source .env  (or use python-dotenv)\n")
        sys.exit(1)
    
    GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY')
    GROQ_API_KEY = os.environ.get('GROQ_API_KEY')
    
    CORS_ORIGINS = os.environ.get('CORS_ORIGINS', 'http://localhost:5000').split(',')
    LOG_LEVEL = os.environ.get('LOG_LEVEL', 'INFO')
    
    DISABLE_AUTH = os.environ.get('DISABLE_AUTH', 'true').lower() == 'true'
    DEBUG = os.environ.get('FLASK_DEBUG', 'False').lower() == 'true'

    # Performance tuning
    DATAFRAME_CACHE_SIZE = int(os.environ.get('DATAFRAME_CACHE_SIZE', 1))
    MAX_ROWS_PER_DATASET = int(os.environ.get('MAX_ROWS_PER_DATASET', 500))

# ------------------------------
#  Flask app initialization
# ------------------------------
app = Flask(__name__)
app.config.from_object(Config)
app.wsgi_app = ProxyFix(app.wsgi_app, x_proto=1, x_host=1)
app.secret_key = Config.SECRET_KEY

CORS(app, origins=Config.CORS_ORIGINS, supports_credentials=True)
Compress(app)                                   # Enable gzip compression

# ------------------------------
#  Rate Limiter
# ------------------------------
class RateLimiter:
    def __init__(self):
        self.requests = defaultdict(list)
    
    def is_allowed(self, key, limit, window_seconds):
        now = datetime.utcnow()
        cutoff = now - timedelta(seconds=window_seconds)
        self.requests[key] = [ts for ts in self.requests[key] if ts > cutoff]
        if len(self.requests[key]) >= limit:
            return False
        self.requests[key].append(now)
        return True

rate_limiter = RateLimiter()

def rate_limit(limit, window=3600, by_ip=True):
    def decorator(f):
        @wraps(f)
        def decorated(*args, **kwargs):
            key = request.remote_addr if by_ip else request.headers.get('Authorization', 'anonymous')
            if not rate_limiter.is_allowed(key, limit, window):
                return jsonify({'error': f'Rate limit exceeded. Max {limit} requests per {window//60} minutes.'}), 429
            return f(*args, **kwargs)
        return decorated
    return decorator

# ------------------------------
#  Authentication (JWT)
# ------------------------------
class AuthManager:
    def __init__(self):
        self.api_keys = {}
    
    def generate_api_key(self, user_id, role='analyst'):
        key = secrets.token_urlsafe(32)
        hashed = hashlib.sha256(key.encode()).hexdigest()
        self.api_keys[hashed] = {'user_id': user_id, 'role': role, 'created_at': datetime.utcnow()}
        return key
    
    def verify_api_key(self, api_key):
        hashed = hashlib.sha256(api_key.encode()).hexdigest()
        return self.api_keys.get(hashed)
    
    def generate_jwt(self, user_id, role):
        payload = {
            'user_id': user_id,
            'role': role,
            'exp': datetime.utcnow() + timedelta(hours=1)
        }
        return jwt.encode(payload, app.config['SECRET_KEY'], algorithm='HS256')

auth_manager = AuthManager()

def require_auth(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if Config.DISABLE_AUTH:
            request.current_user = {'user_id': 'dev_user', 'role': 'admin'}
            return f(*args, **kwargs)
        
        auth_header = request.headers.get('Authorization')
        if not auth_header:
            return jsonify({'error': 'Missing authorization header'}), 401
        try:
            scheme, token = auth_header.split()
            if scheme.lower() != 'bearer':
                return jsonify({'error': 'Invalid auth scheme'}), 401
            payload = jwt.decode(token, app.config['SECRET_KEY'], algorithms=['HS256'])
            request.current_user = payload
        except jwt.ExpiredSignatureError:
            return jsonify({'error': 'Token expired'}), 401
        except jwt.InvalidTokenError:
            return jsonify({'error': 'Invalid token'}), 401
        return f(*args, **kwargs)
    return decorated

def require_role(roles):
    def decorator(f):
        @wraps(f)
        def decorated(*args, **kwargs):
            if Config.DISABLE_AUTH:
                return f(*args, **kwargs)
            if not hasattr(request, 'current_user'):
                return jsonify({'error': 'Authentication required'}), 401
            if request.current_user.get('role') not in roles:
                return jsonify({'error': 'Insufficient permissions'}), 403
            return f(*args, **kwargs)
        return decorated
    return decorator

# ------------------------------
#  Secure Database Pool
# ------------------------------
class SecureDatabase:
    def __init__(self, db_url, min_conn=1, max_conn=10):
        self.pool = SimpleConnectionPool(min_conn, max_conn, db_url)
    
    def get_connection(self):
        return self.pool.getconn()
    
    def put_connection(self, conn):
        self.pool.putconn(conn)
    
    @contextmanager
    def get_cursor(self):
        conn = None
        try:
            conn = self.get_connection()
            with conn.cursor(cursor_factory=RealDictCursor) as cur:
                yield cur
            conn.commit()
        except Exception as e:
            if conn:
                conn.rollback()
            logger.error(f"Database error: {e}")
            raise
        finally:
            if conn:
                self.put_connection(conn)
    
    def execute_query(self, query, params=None):
        with self.get_cursor() as cur:
            cur.execute(query, params)
            if query.strip().upper().startswith(('SELECT', 'WITH')):
                return cur.fetchall()
            return cur.rowcount

db = SecureDatabase(Config.DATABASE_URL)

# ------------------------------
#  Performance indexes (run at startup)
# ------------------------------
def create_performance_indexes():
    """Create indexes on frequently queried columns if they don't exist."""
    index_queries = [
        "CREATE INDEX IF NOT EXISTS idx_fact_orders_order_date ON warehouse.fact_orders(order_date);",
        "CREATE INDEX IF NOT EXISTS idx_fact_orders_customer_id ON warehouse.fact_orders(customer_id);",
        "CREATE INDEX IF NOT EXISTS idx_dim_time_date ON warehouse.dim_time(date);",
        "CREATE INDEX IF NOT EXISTS idx_fact_orders_net_amount ON warehouse.fact_orders(net_amount);"
    ]
    with db.get_cursor() as cur:
        for sql in index_queries:
            try:
                cur.execute(sql)
            except Exception as e:
                logger.warning(f"Index creation failed (may already exist): {e}")

# ------------------------------
#  Input validation helpers
# ------------------------------
def sanitize_output(data):
    if isinstance(data, str):
        return re.sub(r'<[^>]*>', '', data)
    elif isinstance(data, dict):
        return {k: sanitize_output(v) for k, v in data.items()}
    elif isinstance(data, list):
        return [sanitize_output(item) for item in data]
    return data

def validate_nlq_input(question):
    if len(question) > 500:
        return False, "Query too long (max 500 characters)"
    dangerous = [
        r';\s*DROP\s+TABLE', r';\s*DELETE\s+FROM', r';\s*UPDATE\s+.*SET',
        r';\s*INSERT\s+INTO', r'UNION\s+SELECT', r'--\s*$', r'/\*.*\*/'
    ]
    for pattern in dangerous:
        if re.search(pattern, question, re.IGNORECASE):
            return False, "Potentially dangerous query blocked"
    return True, ""

# ------------------------------
#  SQL Utilities (with schema prefix)
# ------------------------------
def add_schema_prefix(sql_query, schema='warehouse'):
    tables = ['fact_orders', 'dim_customers', 'dim_products', 'dim_location',
              'dim_payment', 'dim_status', 'dim_time']
    for table in tables:
        pattern = rf'(?<![\.\w]){table}\b(?!\.)'
        sql_query = re.sub(pattern, f'{schema}.{table}', sql_query, flags=re.IGNORECASE)
    return sql_query

def fix_date_extract(sql_query):
    pattern = r'EXTRACT\(DAY\s+FROM\s+\(CURRENT_DATE\s*-\s*MAX\(([^)]+)\)\)\)'
    sql_query = re.sub(pattern, r'(CURRENT_DATE - MAX(\1))', sql_query, flags=re.IGNORECASE)
    pattern = r'EXTRACT\(DAY\s+FROM\s+\(MAX\(([^)]+)\)\s*-\s*MIN\(([^)]+)\)\)\)'
    sql_query = re.sub(pattern, r'(MAX(\1) - MIN(\2))', sql_query, flags=re.IGNORECASE)
    pattern = r'EXTRACT\(DAY\s+FROM\s+\(([^)]+)\)\)'
    sql_query = re.sub(pattern, r'(\1)', sql_query, flags=re.IGNORECASE)
    return sql_query

def clean_sql(sql_content):
    lines = []
    for line in sql_content.split('\n'):
        if '--' in line:
            line = line[:line.index('--')]
        line = line.strip()
        if line:
            lines.append(line)
    return '\n'.join(lines)

# ------------------------------
#  DataLoader (with reduced cache and row limits)
# ------------------------------
class DataLoader:
    def __init__(self):
        self.sql_folder = self._find_sql_folder()
        self.sql_files = {}
        self._cache = OrderedDict()
        self._cache_maxsize = Config.DATAFRAME_CACHE_SIZE
        self._index_sql_files()

    def _find_sql_folder(self):
        possible = [Path('sql/analytics'), Path('../sql/analytics'), Path('Retail_Analytics/sql/analytics')]
        for p in possible:
            if p.exists():
                return p
        return None

    def _index_sql_files(self):
        if not self.sql_folder:
            return
        for sql_file in self.sql_folder.glob('*.sql'):
            name = re.sub(r'[^\w\-_]', '_', sql_file.stem)
            self.sql_files[name] = sql_file

    def _execute_sql_file(self, sql_path):
        try:
            with open(sql_path, 'r', encoding='utf-8') as f:
                sql_content = f.read()
            sql_content = clean_sql(sql_content)
            if not sql_content:
                return None
            sql_content = add_schema_prefix(sql_content)
            sql_content = fix_date_extract(sql_content)
            # Force a LIMIT to avoid loading huge datasets
            if not re.search(r'\bLIMIT\s+\d+', sql_content, re.IGNORECASE):
                sql_content += f" LIMIT {Config.MAX_ROWS_PER_DATASET}"
            result = db.execute_query(sql_content)
            if result and isinstance(result, list) and len(result) > 0:
                df = pd.DataFrame(result)
                df = self._convert_decimal_to_float(df)
                if len(df) > Config.MAX_ROWS_PER_DATASET:
                    logger.warning(f"Truncating {sql_path.name} from {len(df)} to {Config.MAX_ROWS_PER_DATASET} rows")
                    df = df.head(Config.MAX_ROWS_PER_DATASET)
                # Downcast numeric columns to reduce memory
                for col in df.select_dtypes(include=['float']).columns:
                    df[col] = pd.to_numeric(df[col], downcast='float')
                for col in df.select_dtypes(include=['integer']).columns:
                    df[col] = pd.to_numeric(df[col], downcast='integer')
                return df
            return None
        except Exception as e:
            logger.error(f"Error in {sql_path.name}: {e}")
            return None

    def _convert_decimal_to_float(self, df):
        for col in df.columns:
            try:
                sample = df[col].dropna().iloc[0] if not df[col].dropna().empty else None
                if sample is not None and hasattr(sample, 'as_tuple'):
                    df[col] = df[col].astype(float)
            except Exception:
                pass
        return df

    def get_dataframe(self, sql_key):
        if sql_key in self._cache:
            self._cache.move_to_end(sql_key)
            return self._cache[sql_key]
        sql_path = self.sql_files.get(sql_key)
        if not sql_path:
            return pd.DataFrame()
        df = self._execute_sql_file(sql_path)
        if df is None:
            df = pd.DataFrame()
        self._cache[sql_key] = df
        if len(self._cache) > self._cache_maxsize:
            oldest_key = next(iter(self._cache))
            logger.info(f"Evicting {oldest_key} from cache (size {len(self._cache)-1}/{self._cache_maxsize})")
            del self._cache[oldest_key]
        return df

    @cached_property
    def friendly_data(self):
        mapping = {
            '1a_Customer_Lifetime_Value__CLV_': 'Customer Lifetime Value',
            '1b_Daily_Revenue_Trends': 'Daily Revenue Trends',
            '1c_Monthly_Revenue_Trends': 'Monthly Revenue Trends',
            '1d_Top_Cities_by_Revenue': 'Top Cities by Revenue',
            '2a_Order_Fulfillment_Performance': 'Order Fulfillment Performance',
            '2b_Order_Status_Distribution': 'Order Status Distribution',
            '2c_Revenue_by_Payment_Method': 'Revenue by Payment Method',
            '2d_Repeat_vs_One-Time_Customers': 'Repeat vs One-Time Customers',
            '3a_Cohort_Analysis__Customer_Retention_Over_Time_': 'Cohort Analysis',
            '3b_Customer_Segmentation_using_Revenue___Behavior': 'Customer Segmentation',
            '3c_Order_Value_Distribution___Basket_Analysis': 'Order Value Distribution',
            '3d_Revenue_Contribution_Analysis': 'Revenue Contribution Analysis',
            '3e_Time-to-Purchase_Behavior': 'Time to Purchase Behavior',
            '4_RFM_Segmentation': 'RFM Segmentation',
            '5_Cohort_Retention_Analysis': 'Cohort Retention Analysis',
            '6_Churn_Detection': 'Churn Detection',
            '7a_Revenue_by_Product_Category': 'Revenue by Product Category',
            '7b_Revenue_by_Product_SubCategory': 'Revenue by Product SubCategory',
            '8_Revenue_by_Location': 'Revenue by Location',
            '9_Payment_Method_Analysis': 'Payment Method Analysis',
            '10_Order_Status_Analysis': 'Order Status Analysis'
        }
        fd = {}
        for key, value in mapping.items():
            for res_name in self.sql_files:
                if key in res_name or res_name.startswith(key.split('_')[0]):
                    fd[value] = LazyDataFrame(self, res_name)
                    break
            if value not in fd:
                for res_name in self.sql_files:
                    if value.lower().replace(' ', '_') in res_name.lower():
                        fd[value] = LazyDataFrame(self, res_name)
                        break
            if value not in fd:
                fd[value] = pd.DataFrame()
        return fd

    def to_dict(self, df):
        if df is None or df.empty:
            return []
        return df.replace({np.nan: None}).to_dict(orient='records')


class LazyDataFrame:
    def __init__(self, loader, sql_key):
        self.loader = loader
        self.sql_key = sql_key

    def _get_df(self):
        return self.loader.get_dataframe(self.sql_key)

    def __getitem__(self, key):
        return self._get_df().__getitem__(key)

    def __setitem__(self, key, value):
        self._get_df()[key] = value

    def __getattr__(self, name):
        return getattr(self._get_df(), name)

    def __len__(self):
        return len(self._get_df())

    def __bool__(self):
        return bool(self._get_df())

    def __contains__(self, key):
        return key in self._get_df()

    def to_dict(self, orient='records'):
        return self._get_df().to_dict(orient=orient)

    def copy(self):
        return self._get_df().copy()

loader = DataLoader()
friendly_data = loader.friendly_data

def get_dataset(name):
    return loader.to_dict(friendly_data.get(name, pd.DataFrame()))

# ------------------------------
#  Multi-Provider AI Setup
# ------------------------------
try:
    from google import genai
    from google.genai import types
    GENAI_AVAILABLE = True
except ImportError:
    GENAI_AVAILABLE = False
    logger.warning("google-genai not installed. Gemini disabled.")

try:
    from groq import Groq
    GROQ_AVAILABLE = True
except ImportError:
    GROQ_AVAILABLE = False
    logger.warning("groq not installed. Groq disabled.")

genai_client = None
groq_client = None

if Config.GEMINI_API_KEY and GENAI_AVAILABLE:
    genai_client = genai.Client(api_key=Config.GEMINI_API_KEY)
    logger.info("✅ Gemini AI ready.")
if Config.GROQ_API_KEY and GROQ_AVAILABLE:
    try:
        groq_client = Groq(api_key=Config.GROQ_API_KEY)
        logger.info("✅ Groq AI ready.")
    except Exception as e:
        logger.error(f"Failed to initialise Groq client: {e}")
        groq_client = None

def call_ai_provider(prompt):
    if genai_client:
        try:
            response = genai_client.models.generate_content(
                model="gemini-2.5-flash",
                contents=prompt,
                config=types.GenerateContentConfig(temperature=0.7, max_output_tokens=8192, top_p=0.95)
            )
            logger.info("✅ AI response from Gemini.")
            return response.text
        except Exception as e:
            logger.error(f"Gemini error: {e}")
    if groq_client:
        try:
            completion = groq_client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.7,
                max_tokens=8192,
                top_p=0.95
            )
            logger.info("✅ AI response from Groq.")
            return completion.choices[0].message.content
        except Exception as e:
            logger.error(f"Groq error: {e}")
    logger.warning("No AI provider available. Falling back to local analysis.")
    return None

# ------------------------------
#  API Endpoints (with pagination and aggregation)
# ------------------------------
@app.route('/')
def index():
    return render_template('index.html')

@app.route('/api/login', methods=['POST'])
def login():
    data = request.get_json()
    api_key = data.get('api_key') if data else None
    if not api_key:
        return jsonify({'error': 'API key required'}), 401
    user = auth_manager.verify_api_key(api_key)
    if not user:
        return jsonify({'error': 'Invalid API key'}), 401
    token = auth_manager.generate_jwt(user['user_id'], user['role'])
    return jsonify({'token': token, 'user': user})

@app.route('/api/datasets')
@require_auth
def list_datasets():
    return jsonify(list(friendly_data.keys()))

@app.route('/api/raw/<path:name>')
@require_auth
def raw_dataset(name):
    limit = request.args.get('limit', default=100, type=int)
    offset = request.args.get('offset', default=0, type=int)
    df = friendly_data.get(name)
    if df is None or df.empty:
        return jsonify([])
    sliced = df.iloc[offset:offset+limit]
    return jsonify(sanitize_output(loader.to_dict(sliced)))

@app.route('/api/date_range')
@require_auth
def date_range():
    df = friendly_data.get('Daily Revenue Trends')
    if df is not None and 'order_day' in df.columns:
        dates = pd.to_datetime(df['order_day'])
        return jsonify({'min_date': dates.min().strftime('%Y-%m-%d'), 'max_date': dates.max().strftime('%Y-%m-%d')})
    return jsonify({'min_date': None, 'max_date': None})

@app.route('/api/value_range')
@require_auth
def value_range():
    df = friendly_data.get('Order Value Distribution')
    if df is not None and len(df) > 0 and 'min_order_value' in df.columns:
        return jsonify({'min_value': float(df['min_order_value'].iloc[0]), 'max_value': float(df['max_order_value'].iloc[0])})
    return jsonify({'min_value': 0, 'max_value': 10000})

@app.route('/api/kpis')
@require_auth
def kpis():
    """Direct SQL aggregation – no loading of full DataFrames."""
    with db.get_cursor() as cur:
        cur.execute("SELECT COALESCE(SUM(net_amount), 0) FROM warehouse.fact_orders")
        total_revenue = cur.fetchone()['coalesce']
        cur.execute("SELECT COUNT(DISTINCT order_id) FROM warehouse.fact_orders")
        total_orders = cur.fetchone()['count']
        cur.execute("SELECT COUNT(DISTINCT customer_id) FROM warehouse.fact_orders")
        total_customers = cur.fetchone()['count']
    avg_order = total_revenue / total_orders if total_orders else 0
    return jsonify({
        "total_revenue": float(total_revenue),
        "total_orders": total_orders,
        "total_customers": total_customers,
        "avg_order_value": round(avg_order, 2)
    })

@app.route('/api/revenue_trend', methods=['GET'])
@require_auth
def revenue_trend():
    """Aggregated revenue trend (day, week, month) to reduce payload."""
    granularity = request.args.get('granularity', 'month')  # 'day', 'week', 'month'
    if granularity == 'day':
        sql = """
            SELECT dt.date AS period, SUM(fo.net_amount) AS revenue
            FROM warehouse.fact_orders fo
            JOIN warehouse.dim_time dt ON fo.order_date = dt.date
            GROUP BY dt.date
            ORDER BY dt.date
        """
    elif granularity == 'week':
        sql = """
            SELECT DATE_TRUNC('week', dt.date) AS period, SUM(fo.net_amount) AS revenue
            FROM warehouse.fact_orders fo
            JOIN warehouse.dim_time dt ON fo.order_date = dt.date
            GROUP BY DATE_TRUNC('week', dt.date)
            ORDER BY period
        """
    else:  # month
        sql = """
            SELECT DATE_TRUNC('month', dt.date) AS period, SUM(fo.net_amount) AS revenue
            FROM warehouse.fact_orders fo
            JOIN warehouse.dim_time dt ON fo.order_date = dt.date
            GROUP BY DATE_TRUNC('month', dt.date)
            ORDER BY period
        """
    rows = db.execute_query(sql)
    if not rows:
        return jsonify([])
    result = []
    for row in rows:
        period = row['period']
        result.append({
            'period': period.isoformat() if hasattr(period, 'isoformat') else str(period),
            'revenue': float(row['revenue'])
        })
    return jsonify(sanitize_output(result))

@app.route('/api/daily_revenue')
@require_auth
def daily_revenue():
    limit = request.args.get('limit', default=90, type=int)   # last 90 days
    offset = request.args.get('offset', default=0, type=int)
    df = friendly_data.get('Daily Revenue Trends')
    if df is None or df.empty:
        return jsonify([])
    sliced = df.iloc[offset:offset+limit]
    return jsonify(sanitize_output(loader.to_dict(sliced)))

@app.route('/api/monthly_revenue')
@require_auth
def monthly_revenue():
    limit = request.args.get('limit', default=24, type=int)   # last 24 months
    offset = request.args.get('offset', default=0, type=int)
    df = friendly_data.get('Monthly Revenue Trends')
    if df is None or df.empty:
        return jsonify([])
    sliced = df.iloc[offset:offset+limit]
    return jsonify(sanitize_output(loader.to_dict(sliced)))

@app.route('/api/top_cities')
@require_auth
def top_cities():
    limit = request.args.get('limit', default=10, type=int)
    offset = request.args.get('offset', default=0, type=int)
    df = friendly_data.get('Top Cities by Revenue')
    if df is not None:
        sliced = df.iloc[offset:offset+limit]
        return jsonify(sanitize_output(loader.to_dict(sliced)))
    return jsonify([])

@app.route('/api/revenue_by_category')
@require_auth
def revenue_by_category():
    limit = request.args.get('limit', default=100, type=int)
    offset = request.args.get('offset', default=0, type=int)
    df = friendly_data.get('Revenue by Product Category')
    if df is not None:
        sliced = df.iloc[offset:offset+limit]
        return jsonify(sanitize_output(loader.to_dict(sliced)))
    return jsonify([])

@app.route('/api/revenue_by_subcategory')
@require_auth
def revenue_by_subcategory():
    limit = request.args.get('limit', default=100, type=int)
    offset = request.args.get('offset', default=0, type=int)
    df = friendly_data.get('Revenue by Product SubCategory')
    if df is not None and not df.empty:
        df = df.copy()
        subcat_col = None
        for col in df.columns:
            if 'subcategory' in col.lower() or 'sub_category' in col.lower():
                subcat_col = col
                break
        rev_col = None
        for col in df.columns:
            if 'revenue' in col.lower() or 'amount' in col.lower():
                rev_col = col
                break
        if subcat_col and rev_col:
            df = df.rename(columns={subcat_col: 'subcategory', rev_col: 'revenue'})
        sliced = df.iloc[offset:offset+limit]
        return jsonify(sanitize_output(loader.to_dict(sliced)))
    return jsonify([])

@app.route('/api/revenue_contribution')
@require_auth
def revenue_contribution():
    limit = request.args.get('limit', default=100, type=int)
    offset = request.args.get('offset', default=0, type=int)
    df = friendly_data.get('Revenue Contribution Analysis')
    if df is not None:
        df = df.sort_values('total_revenue', ascending=False)
        df['total_revenue'] = df['total_revenue'].astype(float)
        df['cumulative_percentage'] = (df['total_revenue'].cumsum() / df['total_revenue'].sum()) * 100
        sliced = df.iloc[offset:offset+limit]
        return jsonify(sanitize_output(loader.to_dict(sliced)))
    return jsonify([])

@app.route('/api/order_value_distribution')
@require_auth
def order_value_distribution():
    limit = request.args.get('limit', default=100, type=int)
    offset = request.args.get('offset', default=0, type=int)
    df = friendly_data.get('Order Value Distribution')
    if df is not None:
        sliced = df.iloc[offset:offset+limit]
        return jsonify(sanitize_output(loader.to_dict(sliced)))
    return jsonify([])

@app.route('/api/customer_clv')
@require_auth
def customer_clv():
    df = friendly_data.get('Customer Lifetime Value')
    if df is not None and len(df) > 0:
        if 'category' in df.columns:
            highest = df[df['category'] == 'Highest'].head(5)
            lowest = df[df['category'] == 'Lowest'].head(5)
        else:
            amount_col = 'total_net_amount' if 'total_net_amount' in df.columns else 'total_revenue'
            df[amount_col] = df[amount_col].astype(float)
            sorted_df = df.sort_values(amount_col, ascending=False)
            highest = sorted_df.head(5)
            lowest = sorted_df.tail(5)
        return jsonify({'highest': sanitize_output(loader.to_dict(highest)), 'lowest': sanitize_output(loader.to_dict(lowest))})
    return jsonify({'highest': [], 'lowest': []})

@app.route('/api/repeat_vs_onetime')
@require_auth
def repeat_vs_onetime():
    return jsonify(sanitize_output(get_dataset('Repeat vs One-Time Customers')))

@app.route('/api/customer_segmentation')
@require_auth
def customer_segmentation():
    limit = request.args.get('limit', default=100, type=int)
    offset = request.args.get('offset', default=0, type=int)
    df = friendly_data.get('Customer Segmentation')
    if df is not None:
        df = df.copy()
        if 'total_revenue' in df.columns:
            df['total_revenue'] = df['total_revenue'].astype(float)
        try:
            df['segment'] = pd.qcut(df['total_revenue'], q=4, labels=['Bronze', 'Silver', 'Gold', 'Platinum'])
        except Exception:
            revenue_median = df['total_revenue'].median()
            revenue_high = df['total_revenue'].quantile(0.75)
            revenue_low = df['total_revenue'].quantile(0.25)
            df['segment'] = 'Bronze'
            df.loc[df['total_revenue'] > revenue_low, 'segment'] = 'Silver'
            df.loc[df['total_revenue'] > revenue_median, 'segment'] = 'Gold'
            df.loc[df['total_revenue'] > revenue_high, 'segment'] = 'Platinum'
        sliced = df.iloc[offset:offset+limit]
        return jsonify(sanitize_output(loader.to_dict(sliced)))
    return jsonify([])

@app.route('/api/churn_detection')
@require_auth
def churn_detection():
    limit = request.args.get('limit', default=100, type=int)
    offset = request.args.get('offset', default=0, type=int)
    df = friendly_data.get('Churn Detection')
    if df is not None:
        sliced = df.iloc[offset:offset+limit]
        return jsonify(sanitize_output(loader.to_dict(sliced)))
    return jsonify([])

@app.route('/api/order_status')
@require_auth
def order_status():
    return jsonify(sanitize_output(get_dataset('Order Status Distribution')))

@app.route('/api/payment_methods')
@require_auth
def payment_methods():
    return jsonify(sanitize_output(get_dataset('Payment Method Analysis')))

@app.route('/api/fulfillment_performance')
@require_auth
def fulfillment_performance():
    limit = request.args.get('limit', default=100, type=int)
    offset = request.args.get('offset', default=0, type=int)
    df = friendly_data.get('Order Fulfillment Performance')
    if df is not None:
        sliced = df.iloc[offset:offset+limit]
        return jsonify(sanitize_output(loader.to_dict(sliced)))
    return jsonify([])

@app.route('/api/time_to_purchase')
@require_auth
def time_to_purchase():
    limit = request.args.get('limit', default=100, type=int)
    offset = request.args.get('offset', default=0, type=int)
    df = friendly_data.get('Time to Purchase Behavior')
    if df is not None:
        df_filtered = df[df['days_between_orders'] > 7]
        sliced = df_filtered.iloc[offset:offset+limit]
        return jsonify(sanitize_output(loader.to_dict(sliced)))
    return jsonify([])

@app.route('/api/rfm_segmentation')
@require_auth
def rfm_segmentation():
    limit = request.args.get('limit', default=500, type=int)
    offset = request.args.get('offset', default=0, type=int)
    df = friendly_data.get('RFM Segmentation')
    if df is not None:
        sliced = df.iloc[offset:offset+limit]
        return jsonify(sanitize_output(loader.to_dict(sliced)))
    return jsonify([])

@app.route('/api/cohort_retention')
@require_auth
def cohort_retention():
    limit = request.args.get('limit', default=200, type=int)
    offset = request.args.get('offset', default=0, type=int)
    df = friendly_data.get('Cohort Retention Analysis')
    if df is not None:
        df = df[df['month_number'] > 0]
        sliced = df.iloc[offset:offset+limit]
        return jsonify(sanitize_output(loader.to_dict(sliced)))
    return jsonify([])

@app.route('/api/revenue_by_location')
@require_auth
def revenue_by_location():
    limit = request.args.get('limit', default=100, type=int)
    offset = request.args.get('offset', default=0, type=int)
    df = friendly_data.get('Revenue by Location')
    if df is not None:
        sliced = df.iloc[offset:offset+limit]
        return jsonify(sanitize_output(loader.to_dict(sliced)))
    return jsonify([])

@app.route('/api/revenue_anomalies')
@require_auth
def revenue_anomalies():
    df = friendly_data.get('Daily Revenue Trends')
    if df is None or df.empty:
        return jsonify([])
    date_col = None
    for col in df.columns:
        if col in ['order_day', 'order_date', 'day', 'date', 'transaction_date'] or 'date' in col.lower():
            date_col = col
            break
    rev_col = None
    for col in df.columns:
        if col in ['total_amount', 'revenue', 'amount', 'total_revenue', 'sales'] or 'amount' in col.lower() or 'revenue' in col.lower():
            rev_col = col
            break
    if not date_col or not rev_col:
        return jsonify([])
    df[date_col] = pd.to_datetime(df[date_col])
    if df.groupby(date_col).size().max() > 1:
        df = df.groupby(date_col)[rev_col].sum().reset_index()
    else:
        df = df[[date_col, rev_col]].copy()
    df = df.sort_values(date_col)
    df['pct_change'] = df[rev_col].pct_change() * 100
    anomalies = df[df['pct_change'] < -20][[date_col, rev_col, 'pct_change']]
    anomalies = anomalies.rename(columns={date_col: 'date', rev_col: 'revenue', 'pct_change': 'drop_percent'})
    anomalies['date'] = anomalies['date'].dt.strftime('%Y-%m-%d')
    return jsonify(sanitize_output(loader.to_dict(anomalies)))

@app.route('/api/high_risk_customers')
@require_auth
def high_risk_customers():
    limit = request.args.get('limit', default=20, type=int)
    offset = request.args.get('offset', default=0, type=int)
    rfm_df = friendly_data.get('RFM Segmentation')
    if rfm_df is None or rfm_df.empty:
        return jsonify([])
    required_cols = ['customer_id', 'recency_days', 'frequency', 'monetary']
    for col in required_cols:
        if col not in rfm_df.columns:
            return jsonify([])
    if 'segment' not in rfm_df.columns:
        rec_median = rfm_df['recency_days'].median()
        mon_median = rfm_df['monetary'].median()
        rfm_df['segment'] = 'Others'
        rfm_df.loc[(rfm_df['recency_days'] <= rec_median/2) & (rfm_df['monetary'] >= mon_median*2), 'segment'] = 'Champions'
        rfm_df.loc[(rfm_df['recency_days'] <= rec_median) & (rfm_df['monetary'] >= mon_median), 'segment'] = 'Loyal'
        rfm_df.loc[(rfm_df['recency_days'] > rec_median*1.5) & (rfm_df['monetary'] < mon_median), 'segment'] = 'At Risk'
    monetary_90th = rfm_df['monetary'].quantile(0.9)
    high_risk = rfm_df[(rfm_df['segment'] == 'At Risk') & (rfm_df['monetary'] > monetary_90th)]
    high_risk = high_risk.sort_values('monetary', ascending=False)
    if 'full_name' in high_risk.columns:
        high_risk['full_name'] = high_risk['full_name'].fillna(high_risk['customer_id'].apply(lambda x: f"Customer {x}"))
    else:
        high_risk['full_name'] = high_risk['customer_id'].apply(lambda x: f"Customer {x}")
    result_cols = ['customer_id', 'full_name', 'recency_days', 'frequency', 'monetary', 'segment']
    sliced = high_risk[result_cols].iloc[offset:offset+limit]
    return jsonify(sanitize_output(loader.to_dict(sliced)))

@app.route('/api/aov_by_category')
@require_auth
def aov_by_category():
    df = friendly_data.get('Revenue by Product Category')
    if df is None or df.empty:
        return jsonify([])
    df['aov'] = df['revenue'] / 1000
    return jsonify(sanitize_output(loader.to_dict(df[['category', 'aov']])))

@app.route('/api/frequency_by_category')
@require_auth
def frequency_by_category():
    df = friendly_data.get('Revenue by Product Category')
    if df is None or df.empty:
        return jsonify([])
    df['frequency'] = 1.5
    return jsonify(sanitize_output(loader.to_dict(df[['category', 'frequency']])))

# ------------------------------
#  NLQ Endpoint
# ------------------------------
def get_schema_description():
    schema_name = 'warehouse'
    tables_query = """
        SELECT table_name 
        FROM information_schema.tables 
        WHERE table_schema = %s AND table_type = 'BASE TABLE'
        ORDER BY table_name;
    """
    tables = db.execute_query(tables_query, (schema_name,))
    if not tables:
        return "No tables found in warehouse schema."
    
    description = "Database schema for retail analytics (schema: warehouse):\n\n"
    for tbl in tables:
        table = tbl['table_name']
        cols_query = """
            SELECT column_name, data_type 
            FROM information_schema.columns 
            WHERE table_schema = %s AND table_name = %s
            ORDER BY ordinal_position;
        """
        columns = db.execute_query(cols_query, (schema_name, table))
        if columns:
            col_list = ", ".join([f"{c['column_name']} ({c['data_type']})" for c in columns])
            description += f"- warehouse.{table}: {col_list}\n"
        else:
            description += f"- warehouse.{table}: (no columns found)\n"
    
    description += "\nKey metrics (derived): total_revenue = sum(net_amount), total_orders = count(distinct order_id)\n"
    try:
        date_range_query = "SELECT MIN(date) AS min_date, MAX(date) AS max_date FROM warehouse.dim_time WHERE date IS NOT NULL"
        date_range = db.execute_query(date_range_query)
        if date_range and date_range[0]['min_date'] and date_range[0]['max_date']:
            min_date = date_range[0]['min_date']
            max_date = date_range[0]['max_date']
            description += f"\nIMPORTANT: The data covers dates from {min_date} to {max_date}. If a user asks for a date outside this range, inform them that no data exists for that period.\n"
    except Exception as e:
        logger.warning(f"Could not fetch date range: {e}")
    return description

def generate_sql_from_question(question, previous_error=None):
    schema_desc = get_schema_description()
    if previous_error:
        prompt = f"""You are an expert SQL generator. The previous SQL query failed with this error:

{previous_error}

Please correct the SQL query. Use the schema below.

Schema:
{schema_desc}

User question: {question}

Rules:
- Use schema prefix 'warehouse.'
- Only SELECT statements.
- Return ONLY SQL, no explanation.

Corrected SQL query:
"""
    else:
        prompt = f"""You are an expert SQL generator. Given the schema, convert the question to PostgreSQL SELECT.

Schema:
{schema_desc}

User question: {question}

Rules:
- Use 'warehouse.'
- Only SELECT.
- Return ONLY SQL.

SQL query:
"""
    sql = call_ai_provider(prompt)
    if not sql:
        return None
    sql = re.sub(r'```sql\n?', '', sql)
    sql = re.sub(r'```\n?', '', sql)
    sql = sql.strip()
    if sql.startswith('--'):
        return None
    return sql

nlq_cache = TTLCache(maxsize=50, ttl=3600)

@app.route('/api/ask', methods=['GET'])
@require_auth
@rate_limit(limit=Config.RATELIMIT_NLQ, window=3600, by_ip=True)
def natural_language_query():
    question = request.args.get('q', '').strip()
    if not question:
        return jsonify({"error": "Missing 'q' parameter"}), 400
    
    is_valid, msg = validate_nlq_input(question)
    if not is_valid:
        return jsonify({"error": msg}), 400

    cache_key = hashlib.md5(question.encode()).hexdigest()
    if cache_key in nlq_cache:
        return jsonify({"question": question, "results": nlq_cache[cache_key], "cached": True})

    max_attempts = 2
    last_error = None
    sql = None

    for attempt in range(max_attempts):
        sql = generate_sql_from_question(question, previous_error=last_error)
        if not sql or sql == '-- impossible request':
            return jsonify({"error": "Could not generate SQL", "question": question}), 400

        try:
            test_sql = add_schema_prefix(sql)
            test_sql = fix_date_extract(test_sql)
            if not re.search(r'\bLIMIT\s+\d+', test_sql, re.IGNORECASE):
                test_sql = test_sql.rstrip(';') + " LIMIT 0"
            db.execute_query(test_sql)
            break
        except Exception as e:
            last_error = str(e)
            logger.warning(f"NLQ attempt {attempt+1} failed: {last_error}")
            if attempt == max_attempts - 1:
                return jsonify({"error": f"SQL failed after {max_attempts} attempts", "db_error": last_error}), 500

    try:
        sql = add_schema_prefix(sql)
        sql = fix_date_extract(sql)
        results = db.execute_query(sql)
        serializable = []
        for row in (results or []):
            row_dict = dict(row)
            for k, v in row_dict.items():
                if hasattr(v, 'isoformat'):
                    row_dict[k] = v.isoformat()
            serializable.append(row_dict)
        nlq_cache[cache_key] = serializable
        return jsonify({"question": question, "sql": sql, "results": serializable, "cached": False})
    except Exception as e:
        return jsonify({"error": f"Execution failed: {str(e)}", "sql": sql}), 500

# ------------------------------
#  Simulation Model
# ------------------------------
SIMULATION_COEFFS = {
    "repeat_rate": 1.2,
    "aov": 1.0,
    "churn_rate": 0.8,
    "fulfillment_days": 0.5,
    "last_trained": None
}

def compute_monthly_metrics():
    query = """
    WITH monthly_orders AS (
        SELECT 
            dt.year,
            dt.month,
            fo.customer_id,
            fo.order_id,
            fo.net_amount
        FROM warehouse.fact_orders fo
        JOIN warehouse.dim_time dt ON fo.order_date = dt.date
        WHERE fo.order_date IS NOT NULL
    ),
    customer_monthly AS (
        SELECT 
            year, month,
            customer_id,
            COUNT(DISTINCT order_id) AS order_count,
            SUM(net_amount) AS customer_revenue
        FROM monthly_orders
        GROUP BY year, month, customer_id
    ),
    monthly_agg AS (
        SELECT 
            year, month,
            COUNT(DISTINCT customer_id) AS total_customers,
            COUNT(DISTINCT CASE WHEN order_count > 1 THEN customer_id END) AS repeat_customers,
            SUM(customer_revenue) AS total_revenue,
            COUNT(*) AS total_orders
        FROM customer_monthly
        GROUP BY year, month
    )
    SELECT 
        TO_CHAR(TO_DATE(year || '-' || month || '-01', 'YYYY-MM-DD'), 'YYYY-MM') AS year_month,
        total_revenue,
        (repeat_customers::float / NULLIF(total_customers, 0)) * 100 AS repeat_rate,
        (total_revenue / NULLIF(total_orders, 0)) AS aov
    FROM monthly_agg
    ORDER BY year, month;
    """
    try:
        rows = db.execute_query(query)
        if not rows:
            raise Exception("No monthly data found")
        df = pd.DataFrame(rows)
        df['churn_rate'] = np.maximum(0, 100 - df['repeat_rate'] * 1.5)
        return df
    except Exception as e:
        logger.error(f"Failed to compute monthly metrics: {e}")
        return pd.DataFrame()

def train_simulation_model():
    df = compute_monthly_metrics()
    if df.empty or len(df) < 3:
        logger.warning("Not enough monthly data to train simulation model. Using default coefficients.")
        return

    df = df.sort_values('year_month')
    df['revenue_pct'] = df['total_revenue'].pct_change() * 100
    df['repeat_rate_pct'] = df['repeat_rate'].pct_change() * 100
    df['aov_pct'] = df['aov'].pct_change() * 100
    df['churn_rate_pct'] = df['churn_rate'].pct_change() * 100

    numeric_cols = ['revenue_pct', 'repeat_rate_pct', 'aov_pct', 'churn_rate_pct']
    for col in numeric_cols:
        df[col] = pd.to_numeric(df[col], errors='coerce')
    df_clean = df.replace([np.inf, -np.inf], np.nan).dropna(subset=numeric_cols)

    if len(df_clean) < 2:
        logger.warning("Not enough clean monthly data points for regression. Using default coefficients.")
        return

    X_repeat = df_clean['repeat_rate_pct'].values.reshape(-1, 1).astype(float)
    y_revenue = df_clean['revenue_pct'].values.astype(float)
    if len(X_repeat) >= 2:
        model = LinearRegression()
        model.fit(X_repeat, y_revenue)
        SIMULATION_COEFFS["repeat_rate"] = float(model.coef_[0])
        logger.info(f"Trained repeat_rate coefficient: {SIMULATION_COEFFS['repeat_rate']:.4f}")

    X_aov = df_clean['aov_pct'].values.reshape(-1, 1).astype(float)
    if len(X_aov) >= 2:
        model = LinearRegression()
        model.fit(X_aov, y_revenue)
        SIMULATION_COEFFS["aov"] = float(model.coef_[0])
        logger.info(f"Trained aov coefficient: {SIMULATION_COEFFS['aov']:.4f}")

    X_churn = df_clean['churn_rate_pct'].values.reshape(-1, 1).astype(float)
    if len(X_churn) >= 2:
        model = LinearRegression()
        model.fit(X_churn, y_revenue)
        SIMULATION_COEFFS["churn_rate"] = -float(model.coef_[0])
        logger.info(f"Trained churn_rate coefficient: {SIMULATION_COEFFS['churn_rate']:.4f}")

    SIMULATION_COEFFS["last_trained"] = datetime.now().isoformat()

@app.route('/api/simulate/train', methods=['POST'])
@require_auth
@require_role(['admin'])
def trigger_training():
    try:
        train_simulation_model()
        return jsonify({
            "message": "Simulation model retrained successfully",
            "coefficients": {k: v for k, v in SIMULATION_COEFFS.items() if k != "last_trained"},
            "last_trained": SIMULATION_COEFFS["last_trained"]
        })
    except Exception as e:
        logger.error(f"Training failed: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/simulate', methods=['POST'])
@require_auth
@require_role(['analyst', 'admin'])
@rate_limit(limit=Config.RATELIMIT_SIMULATE, window=3600, by_ip=True)
def simulate():
    data = request.get_json()
    if not data:
        return jsonify({"error": "Invalid request"}), 400
    metric = data.get('metric')
    delta = float(data.get('delta', 0))

    try:
        # Use the new kpis endpoint which is already fast
        total_revenue = kpis().json.get('total_revenue', 0)
        total_orders = kpis().json.get('total_orders', 0)
        aov_current = total_revenue / total_orders if total_orders else 0

        coeff = SIMULATION_COEFFS.get(metric, 1.0)
        if metric == 'churn_rate':
            uplift_pct = delta * coeff
        else:
            uplift_pct = delta * coeff

        estimated_uplift = total_revenue * (uplift_pct / 100)
        new_revenue = total_revenue + estimated_uplift

        extra_info = {}
        if metric == 'aov':
            extra_info = {
                "current_AOV": round(aov_current, 2),
                "new_AOV": round(aov_current * (1 + delta/100), 2)
            }

        return jsonify({
            "metric": metric,
            "delta": delta,
            "estimated_revenue_uplift": round(estimated_uplift, 2),
            "new_total_revenue": round(new_revenue, 2),
            "uplift_percentage": round(uplift_pct, 2),
            "confidence": "data‑driven (regression on monthly data)" if SIMULATION_COEFFS["last_trained"] else "fallback defaults",
            "extra_info": extra_info,
            "coefficient_used": round(coeff, 4)
        })
    except Exception as e:
        logger.error(f"Simulation error: {e}")
        return jsonify({"error": str(e)}), 500

# ------------------------------
#  AI Insights (with fixed repeat rate and order status)
# ------------------------------
PERSONA_TEMPLATES = {
    "conservative_cfo": """
You are a conservative CFO. Focus on cost control, risk, and efficiency.
Be cautious in recommendations. Highlight any financial risks.
Write a concise, actionable business report with these sections:
## Executive Summary (2-3 sentences)
## 1. Key Metrics & Filters
## 2. Deep Dive (Revenue, Retention, Operations)
## 3. Root Causes
## 4. Actionable Recommendations (Short-term, Long-term)
## 5. Expected Business Impact
""",
    "growth_cmo": """
You are an aggressive CMO focused on growth, customer acquisition, and market expansion.
Be optimistic and prioritize recommendations that drive top-line revenue.
Write a concise, actionable business report with these sections:
## Executive Summary (2-3 sentences)
## 1. Key Metrics & Filters
## 2. Deep Dive (Revenue, Retention, Operations)
## 3. Root Causes
## 4. Actionable Recommendations (Short-term, Long-term)
## 5. Expected Business Impact
""",
    "balanced_analyst": """
You are a balanced business analyst. Weigh risks and opportunities equally.
Provide objective, data-driven recommendations.
Write a concise, actionable business report with these sections:
## Executive Summary (2-3 sentences)
## 1. Key Metrics & Filters
## 2. Deep Dive (Revenue, Retention, Operations)
## 3. Root Causes
## 4. Actionable Recommendations (Short-term, Long-term)
## 5. Expected Business Impact
"""
}

feedback_store = {
    "conservative_cfo": {"up": 0, "down": 0},
    "growth_cmo": {"up": 0, "down": 0},
    "balanced_analyst": {"up": 0, "down": 0}
}

def _get_additional_metrics():
    extra = {}
    try:
        extra['fulfillment'] = friendly_data.get('Order Fulfillment Performance')
        extra['payment_methods'] = friendly_data.get('Payment Method Analysis')
        extra['order_status'] = friendly_data.get('Order Status Distribution')
        extra['subcategories'] = friendly_data.get('Revenue by Product SubCategory')
        extra['churn'] = friendly_data.get('Churn Detection')
        extra['cohort_full'] = friendly_data.get('Cohort Retention Analysis')
        extra['rfm_full'] = friendly_data.get('RFM Segmentation')
    except Exception as e:
        logger.warning(f"Could not fetch all extra metrics: {e}")
    return extra

def fix_list_numbering(text):
    pattern = r'(## 5\. Expected Business Impact.*?)(?=\n## |\n---|\Z)'
    match = re.search(pattern, text, re.DOTALL | re.IGNORECASE)
    if not match:
        return text
    section = match.group(1)
    lines = section.split('\n')
    new_lines = []
    counter = 1
    for line in lines:
        if re.match(r'^\s*1\.\s+', line):
            new_line = re.sub(r'^\s*1\.\s+', f'{counter}. ', line)
            new_lines.append(new_line)
            counter += 1
        else:
            new_lines.append(line)
    fixed_section = '\n'.join(new_lines)
    return text.replace(section, fixed_section)

def generate_local_deep_insights_fallback(kpis, filters, daily_revenue, monthly_revenue, top_cities,
                                          revenue_categories, repeat_customers, clv_data, rfm_segments,
                                          cohort_retention, anomalies, high_risk, extra_metrics):
    # Robust repeat rate extraction
    one_time = 0
    repeat_cust = 0
    one_time = repeat_customers.get('one-time', repeat_customers.get('one_time', 0))
    repeat_cust = repeat_customers.get('repeat', repeat_customers.get('repeat_customer', 0))
    if one_time == 0:
        for k, v in repeat_customers.items():
            if 'one' in str(k).lower() and 'time' in str(k).lower():
                one_time = v
                break
    if repeat_cust == 0:
        for k, v in repeat_customers.items():
            if 'repeat' in str(k).lower():
                repeat_cust = v
                break
    total_cust = one_time + repeat_cust
    repeat_rate = (repeat_cust / total_cust * 100) if total_cust else 0
    aov = kpis.get('avg_order_value', 0)
    total_rev = kpis.get('total_revenue', 0)
    
    status_df = extra_metrics.get('order_status')
    status_summary = ""
    if status_df is not None and not status_df.empty:
        total_orders = status_df['order_count'].sum() if 'order_count' in status_df else 0
        status_list = []
        for _, row in status_df.iterrows():
            status = row.get('order_status', 'unknown')
            count = row.get('order_count', 0)
            pct = (count / total_orders * 100) if total_orders else 0
            status_list.append(f"{status}: {count} ({pct:.1f}%)")
        status_summary = ", ".join(status_list[:5])
    else:
        status_summary = "No order status data"
    
    report = []
    report.append("# 📊 Retail Analytics – Deep Business Report (AI Fallback)\n")
    report.append("## Executive Summary\n")
    if repeat_rate < 30:
        report.append(f"⚠️ **Critical**: Only {repeat_rate:.1f}% repeat buyers. Retention is a major risk.\n")
    else:
        report.append(f"✅ **Healthy loyalty**: {repeat_rate:.1f}% repeat rate.\n")
    if anomalies:
        report.append(f"📉 **{len(anomalies)} revenue anomaly days** (>20% drop).\n")
    if high_risk:
        report.append(f"💎 **{len(high_risk)} high‑value customers at risk** of churn.\n")
    report.append(f"Overall: ${total_rev:,.0f} revenue from {kpis.get('total_orders',0):,} orders, AOV ${aov:,.0f}.\n")
    report.append("## 1. Key Metrics & Filters\n")
    date_filter = filters.get('dateRange', {})
    report.append(f"- **Filters**: Date {date_filter.get('min','any')} → {date_filter.get('max','any')}, City {filters.get('selectedCity','any')}, Category {filters.get('selectedCategory','any')}\n")
    report.append(f"- **AOV**: ${aov:,.0f} – " + ("low, consider bundling." if aov < 50 else "healthy.") + "\n")
    report.append(f"- **Repeat Rate**: {repeat_rate:.1f}% – " + ("needs immediate action." if repeat_rate < 30 else "good, aim for 40%+.") + "\n")
    report.append("## 2. Deep Dive\n")
    report.append("### Revenue & Growth\n")
    if top_cities:
        report.append(f"- Top city: {top_cities[0].get('city', 'N/A')}. ")
    if revenue_categories:
        report.append(f"Top category: {revenue_categories[0].get('category', 'N/A')}.\n")
    if anomalies:
        report.append(f"- Anomaly days: {', '.join([a.get('date','')[:10] for a in anomalies[:3]])}.\n")
    report.append("\n### Retention & Loyalty\n")
    report.append(f"- Repeat rate {repeat_rate:.1f}% is " + ("below benchmark." if repeat_rate < 30 else "acceptable.") + "\n")
    if clv_data.get('highest'):
        avg_top_clv = sum(c.get('total_net_amount', c.get('total_revenue', 0)) for c in clv_data['highest']) / len(clv_data['highest'])
        report.append(f"- Top 5 CLV: ${avg_top_clv:,.0f} avg.\n")
    if high_risk:
        report.append(f"- High‑risk VIPs: {len(high_risk)} customers, total ${sum(c.get('monetary',0) for c in high_risk):,.0f} at stake.\n")
    report.append("\n### Operations & Risk\n")
    report.append(f"- Order Status Distribution: {status_summary}\n")
    pay_df = extra_metrics.get('payment_methods')
    top_payment = pay_df.iloc[0].get('payment_method', 'N/A') if pay_df is not None and len(pay_df) > 0 else 'N/A'
    report.append(f"- Top payment method: {top_payment}.\n")
    report.append("\n## 3. Root Causes\n")
    if repeat_rate < 30:
        report.append("- Low repeat rate → weak post‑purchase engagement, no loyalty program.\n")
    if anomalies:
        report.append("- Revenue drops → ended promotions, stockouts, or technical issues.\n")
    if high_risk:
        report.append("- High‑value churn risk → lack of VIP treatment or relevant offers.\n")
    report.append("\n## 4. Actionable Recommendations\n")
    report.append("### Short‑Term (30 days)\n")
    if repeat_rate < 30:
        report.append("1. Launch win‑back email with 15% off for one‑time buyers.\n")
    if aov < 75:
        report.append("2. Create product bundles to increase AOV by 20%.\n")
    if anomalies:
        report.append("3. Set up daily revenue anomaly alerts (Slack/email).\n")
    if high_risk:
        report.append(f"4. Personalised VIP discount codes to top {min(5, len(high_risk))} at‑risk customers.\n")
    if not any([repeat_rate < 30, aov < 75, anomalies, high_risk]):
        report.append("1. Run A/B test on checkout page.\n")
        report.append("2. Introduce referral program.\n")
    report.append("\n### Long‑Term (6‑12 months)\n")
    report.append("- Tiered loyalty program to raise repeat rate to 45%.\n")
    report.append("- Predictive churn model for automated re‑engagement.\n")
    if top_cities:
        report.append(f"- Expand product assortment in {top_cities[0].get('city', 'top city')}.\n")
    report.append("\n## 5. Expected Business Impact\n")
    report.append("1. **Re‑engage VIPs** → recover 20‑30% of lost revenue from high‑risk customers.\n")
    report.append("2. **Loyalty program pilot** → +5‑10% repeat purchase rate within 6 months.\n")
    report.append("3. **Fix fulfillment delays** → reduce cancellations by 1‑2%.\n")
    report.append("\n---\n*Local analysis based on available data.*")
    return "\n".join(report)

def generate_deep_insights_with_persona(kpis, filters, daily_revenue, monthly_revenue, top_cities,
                                         revenue_categories, repeat_customers, clv_data, rfm_segments,
                                         cohort_retention, anomalies, high_risk, extra_metrics,
                                         persona="balanced_analyst"):
    # ----- SAFE REPEAT RATE EXTRACTION -----
    one_time = 0
    repeat_cust = 0
    one_time = repeat_customers.get('one-time', repeat_customers.get('one_time', 
               repeat_customers.get('One-Time', repeat_customers.get('One_Time', 0))))
    repeat_cust = repeat_customers.get('repeat', repeat_customers.get('repeat_customer', 
                                          repeat_customers.get('Repeat', 0)))
    if one_time == 0:
        for k, v in repeat_customers.items():
            k_lower = str(k).lower().replace(' ', '').replace('-', '').replace('_', '')
            if 'one' in k_lower and 'time' in k_lower:
                one_time = v
                break
    if repeat_cust == 0:
        for k, v in repeat_customers.items():
            if 'repeat' in str(k).lower():
                repeat_cust = v
                break
    total_cust = one_time + repeat_cust
    repeat_rate = (repeat_cust / total_cust * 100) if total_cust else 0

    # Order status summary
    status_df = extra_metrics.get('order_status')
    status_summary = ""
    if status_df is not None and not status_df.empty:
        total_orders = status_df['order_count'].sum() if 'order_count' in status_df else 0
        status_list = []
        for _, row in status_df.iterrows():
            status = row.get('order_status', 'unknown')
            count = row.get('order_count', 0)
            pct = (count / total_orders * 100) if total_orders else 0
            status_list.append(f"{status}: {count} ({pct:.1f}%)")
        status_summary = ", ".join(status_list[:5])
    else:
        status_summary = "No order status data available"

    # Build daily and monthly strings
    daily_vals = []
    for d in daily_revenue[-7:]:
        val = d.get('total_amount') or d.get('revenue') or 0
        daily_vals.append(f"${val:,.0f}")
    daily_str = ", ".join(daily_vals) if daily_vals else "no data"
    
    monthly_str = ""
    for m in monthly_revenue[-6:]:
        month = m.get('year_month', 'unknown')
        rev = m.get('total_amount', 0)
        monthly_str += f"{month}: ${rev:,.0f}; "
    
    top_cities_str = ", ".join([f"{c.get('city', 'N/A')} (${c.get('total_revenue',0):,.0f})" for c in top_cities[:3]])
    categories_str = ", ".join([f"{c.get('category', 'N/A')} (${c.get('revenue',0):,.0f})" for c in revenue_categories[:5]])
    
    highest_clv_list = clv_data.get('highest', [])
    avg_top_clv = sum(c.get('total_net_amount', c.get('total_revenue', 0)) for c in highest_clv_list) / max(len(highest_clv_list), 1)
    
    rfm_sample = rfm_segments[:3] if rfm_segments else []
    rfm_str = ", ".join([str(s.get('segment', s.get('rfm_segment', 'unknown'))) for s in rfm_sample])
    
    anomaly_dates = [a.get('date', '')[:10] for a in anomalies[:3]]
    anomalies_str = f"{len(anomalies)} days, e.g. {', '.join(anomaly_dates)}" if anomalies else "none"
    
    high_risk_total = sum(c.get('monetary', 0) for c in high_risk[:5])
    high_risk_str = f"{len(high_risk)} customers, total at-risk ${high_risk_total:,.0f}" if high_risk else "none"
    
    cohort_str = ""
    for c in cohort_retention[:3]:
        cohort_str += f"{c.get('cohort_month', '')} month {c.get('month_number',0)}: {c.get('retention_rate',0)*100:.1f}%; "
    
    churn_df = extra_metrics.get('churn')
    churn_rate_val = churn_df['churn_rate'].iloc[0] if churn_df is not None and 'churn_rate' in churn_df.columns else None
    
    pay_df = extra_metrics.get('payment_methods')
    top_payment = pay_df.iloc[0].get('payment_method', 'N/A') if pay_df is not None and len(pay_df) > 0 else 'N/A'

    context = f"""
KPIs: Revenue ${kpis.get('total_revenue',0):,.0f}, Orders {kpis.get('total_orders',0):,}, Customers {kpis.get('total_customers',0):,}, AOV ${kpis.get('avg_order_value',0):,.0f}
Filters: Date {filters.get('dateRange',{}).get('min','any')} -> {filters.get('dateRange',{}).get('max','any')}, City {filters.get('selectedCity','any')}, Category {filters.get('selectedCategory','any')}
Daily Revenue (last 7 days): {daily_str}
Monthly Revenue (last 6 months): {monthly_str}
Top Cities: {top_cities_str}
Top Categories: {categories_str}
Repeat Rate: {repeat_rate:.1f}%
Avg CLV Top 5: ${avg_top_clv:,.0f}
RFM Segments: {rfm_str}
Cohort Retention (sample): {cohort_str}
Revenue Anomalies: {anomalies_str}
High Risk VIPs: {high_risk_str}
Order Status Distribution: {status_summary}
Top Payment Method: {top_payment}
Churn Rate: {churn_rate_val if churn_rate_val is not None else 'N/A'}%
"""
    base_prompt = PERSONA_TEMPLATES.get(persona, PERSONA_TEMPLATES["balanced_analyst"])
    full_prompt = f"""{base_prompt}

--- DATA ---
{context}
"""
    insights = call_ai_provider(full_prompt)
    if insights:
        required = ["Executive Summary", "Key Metrics", "Deep Dive", "Root Causes", "Actionable Recommendations", "Expected Business Impact"]
        if not any(sec in insights for sec in required):
            insights += "\n\n---\n*Note: Some sections may be abbreviated due to token limits.*"
        return insights + "\n\n---\n*Powered by AI*"
    else:
        return generate_local_deep_insights_fallback(
            kpis, filters, daily_revenue, monthly_revenue, top_cities,
            revenue_categories, repeat_customers, clv_data, rfm_segments,
            cohort_retention, anomalies, high_risk, extra_metrics
        )

ai_insights_cache = TTLCache(maxsize=100, ttl=21600)

@app.route('/api/ai_insights', methods=['POST'])
@require_auth
@rate_limit(limit=Config.RATELIMIT_AI, window=3600, by_ip=True)
def ai_insights():
    data = request.get_json()
    if not data:
        return jsonify({'error': 'No data provided'}), 400
    
    filters = data.get('filters', {})
    kpis_data = data.get('kpis', {})
    daily_revenue = data.get('daily_revenue', [])
    monthly_revenue = data.get('monthly_revenue', [])
    top_cities = data.get('top_cities', [])
    revenue_categories = data.get('revenue_categories', [])
    repeat_customers = data.get('repeat_customers', {})
    clv_data = data.get('clv_data', {'highest': [], 'lowest': []})
    rfm_segments = data.get('rfm_segments', [])
    cohort_retention = data.get('cohort_retention', [])
    anomalies = data.get('anomalies', [])
    high_risk = data.get('high_risk_customers', [])
    persona = data.get('persona', 'balanced_analyst')
    
    cache_key = hashlib.md5(json.dumps({
        "filters": filters,
        "kpis": kpis_data,
        "daily_tail": [d.get('total_amount') for d in daily_revenue[-7:]],
        "monthly_tail": [m.get('total_amount') for m in monthly_revenue[-6:]],
        "top_cities_ids": [c.get('city') for c in top_cities[:3]],
        "persona": persona,
        "repeat_rate_hash": repeat_customers.get('repeat', 0)
    }, sort_keys=True).encode()).hexdigest()
    
    if cache_key in ai_insights_cache:
        insights = ai_insights_cache[cache_key]
    else:
        extra_metrics = _get_additional_metrics()
        insights = generate_deep_insights_with_persona(
            kpis=kpis_data, filters=filters, daily_revenue=daily_revenue,
            monthly_revenue=monthly_revenue, top_cities=top_cities,
            revenue_categories=revenue_categories, repeat_customers=repeat_customers,
            clv_data=clv_data, rfm_segments=rfm_segments, cohort_retention=cohort_retention,
            anomalies=anomalies, high_risk=high_risk, extra_metrics=extra_metrics,
            persona=persona
        )
        ai_insights_cache[cache_key] = insights
    
    insights = fix_list_numbering(insights)
    return jsonify({"insights": insights, "persona": persona})

# ------------------------------
#  Alert, Export, Feedback
# ------------------------------
def send_slack_alert(message, webhook_url=None):
    logger.info(f"SIMULATED SLACK ALERT: {message}")

def create_jira_card(summary, description, jira_config=None):
    logger.info(f"SIMULATED JIRA CARD: {summary}\n{description}")

def check_and_trigger_alerts(kpis, anomalies, repeat_customers, extra_metrics):
    alerts = []
    if anomalies and len(anomalies) > 0:
        alert_msg = f"⚠️ Revenue Anomaly: {len(anomalies)} days with >20% drop."
        alerts.append(alert_msg)
        send_slack_alert(alert_msg)
        create_jira_card("Revenue Anomaly Alert", alert_msg)
    one_time = repeat_customers.get('one-time', 0)
    repeat_cust = repeat_customers.get('repeat', 0)
    total_cust = one_time + repeat_cust
    repeat_rate = (repeat_cust / total_cust * 100) if total_cust else 0
    if repeat_rate < 30:
        alert_msg = f"⚠️ Low Repeat Rate: {repeat_rate:.1f}% (below 30% threshold)"
        alerts.append(alert_msg)
        send_slack_alert(alert_msg)
        create_jira_card("Low Repeat Rate", alert_msg)
    high_risk = extra_metrics.get('high_risk_vips', [])
    if high_risk and len(high_risk) > 5:
        alert_msg = f"⚠️ {len(high_risk)} high-value customers at risk of churn"
        alerts.append(alert_msg)
        send_slack_alert(alert_msg)
        create_jira_card("VIP Churn Risk", alert_msg)
    return alerts

@app.route('/api/check_anomalies', methods=['POST'])
@require_auth
def check_anomalies():
    try:
        kpis_data = kpis().json
        daily_df = friendly_data.get('Daily Revenue Trends')
        anomalies_list = []
        if daily_df is not None and not daily_df.empty:
            rev_col = 'total_amount'
            df_sorted = daily_df.sort_values('order_day')
            df_sorted['pct_change'] = df_sorted[rev_col].pct_change() * 100
            anomalies_df = df_sorted[df_sorted['pct_change'] < -20]
            anomalies_list = anomalies_df[['order_day', rev_col, 'pct_change']].rename(columns={'order_day': 'date', rev_col: 'revenue', 'pct_change': 'drop_percent'}).to_dict(orient='records')
        repeat_data = get_dataset('Repeat vs One-Time Customers')
        repeat_dict = {r['customer_type']: r['customer_count'] for r in repeat_data} if repeat_data else {}
        extra_metrics = _get_additional_metrics()
        extra_metrics['high_risk_vips'] = get_dataset('high_risk_customers')
        alerts = check_and_trigger_alerts(kpis_data, anomalies_list, repeat_dict, extra_metrics)
        return jsonify({"alerts": alerts, "triggered": len(alerts) > 0})
    except Exception as e:
        logger.error(f"Alert check failed: {e}")
        return jsonify({"error": str(e)}), 500

try:
    from weasyprint import HTML
    WEASYPRINT_AVAILABLE = True
except ImportError:
    WEASYPRINT_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

def generate_pdf_from_html(html_content):
    if not WEASYPRINT_AVAILABLE:
        raise ImportError("weasyprint not installed")
    return HTML(string=html_content).write_pdf()

def generate_powerpoint(insights_text):
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx not installed")
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "AI Business Report"
    subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = bullet_slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Executive Summary & Insights"
    tf = body_shape.text_frame
    tf.text = insights_text[:1000]
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        tmp.seek(0)
        pptx_bytes = tmp.read()
    os.unlink(tmp.name)
    return pptx_bytes

@app.route('/api/export', methods=['POST'])
@require_auth
def export_report():
    data = request.get_json()
    format_type = data.get('format', 'pdf')
    insights_html = data.get('insights_html', '<h1>No insights available</h1>')
    try:
        if format_type == 'pdf':
            pdf_bytes = generate_pdf_from_html(insights_html)
            return send_file(
                io.BytesIO(pdf_bytes),
                mimetype='application/pdf',
                as_attachment=True,
                download_name=f'ai_insights_{datetime.now().strftime("%Y%m%d_%H%M%S")}.pdf'
            )
        elif format_type == 'pptx':
            plain_text = re.sub(r'<[^>]+>', '', insights_html)
            pptx_bytes = generate_powerpoint(plain_text)
            return send_file(
                io.BytesIO(pptx_bytes),
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                as_attachment=True,
                download_name=f'ai_insights_{datetime.now().strftime("%Y%m%d_%H%M%S")}.pptx'
            )
        else:
            return jsonify({"error": "Unsupported format"}), 400
    except Exception as e:
        logger.error(f"Export failed: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/feedback', methods=['POST'])
@require_auth
def submit_feedback():
    data = request.get_json()
    persona = data.get('persona')
    feedback = data.get('feedback')
    if persona not in feedback_store:
        return jsonify({"error": "Invalid persona"}), 400
    if feedback not in ['up', 'down']:
        return jsonify({"error": "Feedback must be 'up' or 'down'"}), 400
    feedback_store[persona][feedback] += 1
    logger.info(f"Feedback for {persona}: {feedback} (now up={feedback_store[persona]['up']}, down={feedback_store[persona]['down']})")
    return jsonify({"message": "Thank you for your feedback!", "stats": feedback_store[persona]})

@app.route('/api/feedback/stats', methods=['GET'])
@require_auth
def get_feedback_stats():
    return jsonify(feedback_store)

# ------------------------------
#  Security Headers
# ------------------------------
@app.after_request
def add_security_headers(response):
    response.headers['X-Content-Type-Options'] = 'nosniff'
    response.headers['X-Frame-Options'] = 'DENY'
    response.headers['X-XSS-Protection'] = '1; mode=block'
    response.headers['Strict-Transport-Security'] = 'max-age=31536000; includeSubDomains'
    response.headers['Cache-Control'] = 'no-store, max-age=0'
    return response

@app.route('/health', methods=['GET'])
def health():
    return jsonify({'status': 'healthy', 'timestamp': datetime.utcnow().isoformat()})

# ------------------------------
#  Run app
# ------------------------------
if __name__ == '__main__':
    create_performance_indexes()
    with app.app_context():
        train_simulation_model()
    app.run(host='0.0.0.0', port=5001, debug=Config.DEBUG)